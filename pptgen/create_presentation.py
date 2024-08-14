"""Construct Powerpoint."""

import io
from typing import Any, List

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from pptgen.colors import (
    apply_background_gradient,
    apply_text_formatting,
)
from pptgen.model.powerpoint.common import BulletPoints
from pptgen.model.powerpoint.content_slide import ContentSlide
from pptgen.model.powerpoint.image_slide import ImageSlide
from pptgen.model.powerpoint.title_slide import TitleSlide


def apply_paragraph_formatting(paragraph, slide_model, color_scheme, bullet_point=None):
    """Apply formatting to a paragraph."""
    if bullet_point:
        font_size = bullet_point.font_size
        font_name = bullet_point.font_name
    else:
        font_size = slide_model.content_font_size
        font_name = slide_model.content_font_name

    apply_text_formatting(
        paragraph,
        color_scheme.content_color,
        font_size,
        font_name,
        alignment=PP_ALIGN.LEFT,
    )


def add_title_slide(prs, slide_model, color_scheme):
    """Add a title slide to the presentation."""
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slide_model.title
    subtitle.text = slide_model.subtitle

    # Apply text properties
    apply_text_formatting(
        title.text_frame.paragraphs[0],
        color_scheme.title_color,
        slide_model.title_font_size,
        slide_model.title_font_name,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    apply_text_formatting(
        subtitle.text_frame.paragraphs[0],
        color_scheme.subtitle_color,
        slide_model.subtitle_font_size,
        slide_model.subtitle_font_name,
        alignment=PP_ALIGN.CENTER,
    )

    # Set background gradient
    apply_background_gradient(slide, color_scheme.background_gradient)

    return slide


def add_content_slide(prs, slide_model, color_scheme):
    """Add a content slide to the presentation."""
    layout = prs.slide_layouts[1]  # Content with Caption layout
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = slide_model.title

    # Apply text properties to title
    apply_text_formatting(
        title.text_frame.paragraphs[0],
        color_scheme.title_color,
        slide_model.title_font_size,
        slide_model.title_font_name,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    # Handle content based on its type
    if isinstance(slide_model.content, str):
        content.text = slide_model.content
        for paragraph in content.text_frame.paragraphs:
            apply_paragraph_formatting(paragraph, slide_model, color_scheme)
    elif isinstance(slide_model.content, BulletPoints):
        text_frame = content.text_frame
        text_frame.clear()  # Clear any existing text
        for bullet in slide_model.content.bullet_points:
            p = text_frame.add_paragraph()
            p.text = bullet.text
            p.level = 0  # Set to 0 for top-level bullets
            apply_paragraph_formatting(p, slide_model, color_scheme, bullet)

    # Set background gradient
    apply_background_gradient(slide, color_scheme.background_gradient)

    return slide


def add_image_slide(prs, slide_model, color_scheme):
    """Add an image slide to the presentation."""
    layout = prs.slide_layouts[5]  # Picture with Caption layout
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = slide_model.title

    # Apply text properties to title
    apply_text_formatting(
        title.text_frame.paragraphs[0],
        color_scheme.title_color,
        slide_model.title_font_size,
        slide_model.title_font_name,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Add image
    left = Inches(slide_model.margin_left)
    top = Inches(slide_model.image_margin_top)
    width = Inches(slide_model.image_width)
    height = Inches(slide_model.image_height)
    slide.shapes.add_picture(slide_model.image_path, left, top, width, height)

    # Set background gradient
    apply_background_gradient(slide, color_scheme.background_gradient)

    return slide


def create_presentation(slide_models: List[Any], color_scheme) -> io.BytesIO:
    """Create a complete presentation based on slide models and color scheme."""
    prs = Presentation()

    for slide in slide_models:
        if isinstance(slide, TitleSlide):
            add_title_slide(prs, slide, color_scheme)
        elif isinstance(slide, ContentSlide):
            add_content_slide(prs, slide, color_scheme)
        elif isinstance(slide, ImageSlide):
            add_image_slide(prs, slide, color_scheme)

    # Save to a BytesIO object
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)

    return pptx_file
