"""Generate the dataframe meta data for the given dataframe."""

from io import BytesIO
from typing import List

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation, presentation
from pptx.util import Inches

from pptgen.colors import apply_background_gradient, apply_text_formatting
from pptgen.model.dataframe_meta import ColumnMeta, ColumnsMeta


def get_dataframe_metadata(df: pd.DataFrame) -> ColumnsMeta:
    """Extract metadata from dataframe and return a ColumnsMeta object."""
    metadata = []
    for col in df.columns:
        col_type = str(df[col].dtype)
        non_null = df[col].count()
        null = df[col].isnull().sum()
        unique = df[col].nunique()

        column_meta = ColumnMeta(
            column=col,
            type=col_type,
            non_null_count=non_null,
            null_count=null,
            unique_values=unique,
        )
        metadata.append(column_meta)

    return ColumnsMeta(columns=metadata)


def create_consolidated_view(
    columns_meta: ColumnsMeta, prs: Presentation, color_scheme
) -> List[presentation.Slides]:
    """Create a consolidated view for dataframes with more than 10 columns."""
    slides = []
    chunk_size = 8  # Changed from 50 to 8

    for i in range(0, len(columns_meta.columns), chunk_size):
        chunk = columns_meta.columns[i : i + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Table slide layout

        # Add title
        title = slide.shapes.title
        title.text = f"DataFrame Metadata (Columns {i+1}-{i+len(chunk)})"
        apply_text_formatting(
            title.text_frame.paragraphs[0],
            color_scheme.title_color,
            24,
            "Arial",
            bold=True,
        )

        # Add table
        rows, cols = len(chunk) + 1, 3  # +1 for header
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set header
        headers = ["Column", "Type", "Non-Null Count"]
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            apply_text_formatting(
                cell.text_frame.paragraphs[0],
                color_scheme.subtitle_color,
                14,
                "Arial",
                bold=True,
            )

        # Fill data
        for row, col_data in enumerate(chunk, start=1):
            table.cell(row, 0).text = col_data.column
            table.cell(row, 1).text = col_data.type
            table.cell(row, 2).text = str(col_data.non_null_count)

        # Apply background gradient
        apply_background_gradient(slide, color_scheme.background_gradient)

        slides.append(slide)

    return slides


def create_detailed_view(
    columns_meta: ColumnsMeta, prs: Presentation, color_scheme
) -> List[presentation.Slides]:
    """Create a detailed view for dataframes with 10 or fewer columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Table slide layout

    # Add title
    title = slide.shapes.title
    title.text = "DataFrame Metadata"
    apply_text_formatting(
        title.text_frame.paragraphs[0], color_scheme.title_color, 24, "Arial", bold=True
    )

    # Add table
    rows, cols = len(columns_meta.columns) + 1, 5  # +1 for header
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set header
    headers = ["Column", "Type", "Non-Null Count", "Null Count", "Unique Values"]
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        apply_text_formatting(
            cell.text_frame.paragraphs[0],
            color_scheme.subtitle_color,
            14,
            "Arial",
            bold=True,
        )

    # Fill data
    for row, col_data in enumerate(columns_meta.columns, start=1):
        table.cell(row, 0).text = col_data.column
        table.cell(row, 1).text = col_data.type
        table.cell(row, 2).text = str(col_data.non_null_count)
        table.cell(row, 3).text = str(col_data.null_count)
        table.cell(row, 4).text = str(col_data.unique_values)

    # Apply background gradient
    apply_background_gradient(slide, color_scheme.background_gradient)

    return [slide]


def remove_last_monthly_count(monthly_counts: pd.DataFrame) -> pd.DataFrame:
    """Remove the last monthly count."""
    # Check if the last month should be excluded
    if len(monthly_counts) > 1:
        last_month_count = monthly_counts.iloc[-1][0]
        previous_month_count = monthly_counts.iloc[-2][0]
        if last_month_count < 0.5 * previous_month_count:
            monthly_counts = monthly_counts[:-1]  # Exclude the last month

    return monthly_counts


def create_df_monthly_counts(df: pd.DataFrame) -> pd.DataFrame:
    """Create the dataframe for monthly counts."""
    df["date"] = pd.to_datetime(df["FILE_DATE"], format="%Y-%m-%d")
    # Create a cumulative count of filings
    monthly_counts = df.groupby(df["date"].dt.to_period("M")).size().reset_index()
    monthly_counts["date"] = monthly_counts["date"].dt.to_timestamp()

    # Remove the last monthly count if it is less than half of the previous month
    return remove_last_monthly_count(monthly_counts)


def plot_monthly_counts(monthly_counts: pd.DataFrame, company_name: str) -> BytesIO:
    """Plot the monthly counts."""
    plt.figure(figsize=(6, 4))
    plt.plot(monthly_counts["date"], monthly_counts[0], color="#0066CC")
    plt.title(
        f"Monthly Count of Filings for {company_name.upper()}",
        fontsize=14,
        fontweight="bold",
    )
    plt.xlabel("Date", fontsize=10)
    plt.ylabel("Count", fontsize=10)
    plt.grid(True, linestyle="--", alpha=0.7)
    plt.xticks(rotation=45)
    plt.tight_layout()

    # Save the plot to a BytesIO object
    img_bytes = BytesIO()
    plt.savefig(img_bytes, format="png", dpi=300, bbox_inches="tight")
    img_bytes.seek(0)

    return img_bytes


def add_plot_to_slide(
    img_bytes: BytesIO,
    slide: Presentation,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Add a plot to a slide."""
    # Add the graph to the slide
    slide.shapes.add_picture(img_bytes, left, top, width, height)

    plt.close()


def create_overview_slide(
    df: pd.DataFrame, prs: Presentation, color_scheme, company_name: str
) -> presentation.Slides:
    """Create an overview slide with total row count and rows per year."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content layout

    # Add title
    title = slide.shapes.title
    title.text = f"{company_name.upper()} Overview"
    apply_text_formatting(
        title.text_frame.paragraphs[0], color_scheme.title_color, 24, "Arial", bold=True
    )

    # Add total row count
    total_rows = len(df)
    total_rows_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(0.5)
    )
    total_rows_text = total_rows_shape.text_frame.add_paragraph()
    total_rows_text.text = f"Total number of rows: {total_rows:,}"
    apply_text_formatting(
        total_rows_text, color_scheme.content_color, 18, "Arial", bold=True
    )

    # Set dimensions for both table and graph
    left_margin = Inches(0.5)
    top_margin = Inches(2.5)
    width = Inches(4.5)
    height = Inches(4.5)

    # Find year column
    if "FILE_YEAR" in df.columns:
        year_column = "FILE_YEAR"
    else:
        year_columns = (
            df.filter(regex=r"year", case=False)
            .select_dtypes(include=["int64", "float64"])
            .columns
        )
        year_column = year_columns[0] if len(year_columns) > 0 else None

    if year_column:
        # Get value counts for the year column
        year_counts = df[year_column].value_counts().sort_index()

        # Add table
        # rows, cols = len(year_counts) + 1, 2  # +1 for header
        # Add table
        rows = min(len(year_counts) + 1, 20)  # Limit to 19 data rows + header
        cols = 2
        # left = Inches(0.5)
        # top = Inches(2.5)
        # width = Inches(4)
        # height = Inches(0.5 * (rows + 1))
        # height = Inches(4.5)  # Fixed height
        table = slide.shapes.add_table(
            rows, cols, left_margin, top_margin, width, height
        ).table

        # Set header
        headers = ["Year", "Number of Rows"]
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            apply_text_formatting(
                cell.text_frame.paragraphs[0],
                color_scheme.subtitle_color,
                12,
                "Arial",
                bold=True,
            )

        # Fill data
        for row, (year, count) in enumerate(year_counts.items(), start=1):
            if row >= rows:
                break
            table.cell(row, 0).text = str(year)
            table.cell(row, 1).text = f"{count:,}"
            apply_text_formatting(
                table.cell(row, 0).text_frame.paragraphs[0],
                color_scheme.content_color,
                10,
                "Arial",
            )
            apply_text_formatting(
                table.cell(row, 1).text_frame.paragraphs[0],
                color_scheme.content_color,
                10,
                "Arial",
            )
        # Adjust row heights
        for row in table.rows:
            row.height = int(height / rows)
    else:
        # If no year column found, add a message
        no_year_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(0.5)
        )
        no_year_text = no_year_shape.text_frame.add_paragraph()
        no_year_text.text = "No year column found in the DataFrame."
        apply_text_formatting(
            no_year_text, color_scheme.content_color, 14, "Arial", italic=True
        )

    monthly_counts = create_df_monthly_counts(df.copy())

    img_bytes = plot_monthly_counts(monthly_counts, company_name)

    graph_left = Inches(5.5)  # Adjusted to create some space between table and graph
    add_plot_to_slide(img_bytes, slide, graph_left, top_margin, width, height)

    # Apply background gradient
    apply_background_gradient(slide, color_scheme.background_gradient)

    return slide
